from flask import Flask, render_template, request, send_from_directory, send_file, redirect, url_for, jsonify
import os
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader, PdfWriter
import win32com.client
import io
import tempfile

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# 零宽字符编码
ZERO_WIDTH_SPACE = '\u200b'  # 零宽空格
ZERO_WIDTH_NON_JOINER = '\u200c'  # 零宽非连接符
ZERO_WIDTH_JOINER = '\u200d'  # 零宽连接符

# 加密：将二进制数据转换为零宽字符
def encrypt_text(text, secret):
    try:
        if isinstance(text, bytes):
            text = text.decode('utf-8')
        if isinstance(secret, bytes):
            secret = secret.decode('utf-8')
        
        # 将密钥转换为二进制字符串
        binary_secret = ''.join(format(ord(c), '08b') for c in secret)
        
        # 添加固定的开始和结束标记
        marker = '11111111'  # 8个1作为标记
        final_binary = marker + binary_secret + marker
        
        print(f"Binary secret: {binary_secret}")
        print(f"Final binary: {final_binary}")
        print(f"Binary length: {len(final_binary)}")
        
        # 使用零宽字符替换二进制位
        zero_width_secret = ''
        for bit in final_binary:
            if bit == '0':
                zero_width_secret += ZERO_WIDTH_SPACE
            else:
                zero_width_secret += ZERO_WIDTH_JOINER
        
        result = text + zero_width_secret
        print(f"Original text length: {len(text)}")
        print(f"Encrypted text length: {len(result)}")
        print(f"Zero-width characters: {len(zero_width_secret)}")
        return result
    except Exception as e:
        print(f"Encryption error: {str(e)}")
        return text

# 解密：从文本中提取零宽字符编码
def decrypt_text(text):
    try:
        if isinstance(text, bytes):
            text = text.decode('utf-8')
        
        # 提取零宽字符
        zero_width_chars = []
        for c in text:
            if c in [ZERO_WIDTH_SPACE, ZERO_WIDTH_JOINER]:
                zero_width_chars.append(c)
        
        print(f"Found {len(zero_width_chars)} zero-width characters")
        
        if not zero_width_chars:
            print("No zero-width characters found")
            return ""
        
        # 转换回二进制字符串
        binary_secret = ''
        for c in zero_width_chars:
            if c == ZERO_WIDTH_SPACE:
                binary_secret += '0'
            elif c == ZERO_WIDTH_JOINER:
                binary_secret += '1'
        
        print(f"Full binary string: {binary_secret}")
        
        # 查找开始和结束标记
        marker = '11111111'
        start_pos = binary_secret.find(marker)
        if start_pos == -1:
            print("Start marker not found")
            return ""
            
        end_pos = binary_secret.rfind(marker)  # 使用 rfind 找到最后一个标记
        if end_pos == -1 or end_pos == start_pos:
            print("End marker not found")
            return ""
        
        # 提取实际的二进制数据（不包括标记）
        binary_data = binary_secret[start_pos + 8:end_pos]
        print(f"Extracted binary data: {binary_data}")
        print(f"Binary data length: {len(binary_data)}")
        
        # 确保二进制数据长度是8的倍数
        if len(binary_data) % 8 != 0:
            # 尝试补全缺失的位
            missing_bits = 8 - (len(binary_data) % 8)
            binary_data = binary_data + '0' * missing_bits
            print(f"Padded binary data: {binary_data}")
            print(f"New binary data length: {len(binary_data)}")
        
        # 转换为字符
        secret = ''
        for i in range(0, len(binary_data), 8):
            byte = binary_data[i:i+8]
            try:
                char_code = int(byte, 2)
                secret += chr(char_code)
                print(f"Converted byte {byte} to character {chr(char_code)}")
            except ValueError as e:
                print(f"Error converting binary to character: {str(e)}")
                continue
        
        print(f"Decrypted secret: {secret}")
        return secret
    except Exception as e:
        print(f"Decryption error: {str(e)}")
        return ""

@app.route('/')
def index():
    return redirect(url_for('text_page'))

@app.route('/text')
def text_page():
    return render_template('text.html')

@app.route('/file')
def file_page():
    return render_template('file.html')

@app.route('/encrypt', methods=['POST'])
def encrypt():
    text = request.form['text']
    secret = request.form['secret']
    result = encrypt_text(text, secret)
    return jsonify({'result': result, 'action': '加密'})

@app.route('/decrypt', methods=['POST'])
def decrypt():
    text = request.form['text']
    result = decrypt_text(text)
    return jsonify({'result': result, 'action': '解密'})

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return render_template('index.html', result="未找到文件")
    file = request.files['file']
    if file.filename == '':
        return render_template('index.html', result="未选择文件")
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(filepath)
    with open(filepath, 'r', encoding='utf-8') as f:
        text = f.read()
    result = decrypt_text(text)
    os.remove(filepath)
    return render_template('index.html', result=result, action='文件解密')

# 添加新的文件处理函数
def process_docx(file_path, mode='encrypt', secret=''):
    doc = Document(file_path)
    result = []
    
    try:
        if mode == 'encrypt':
            modified = False
            for para in doc.paragraphs:
                if para.text.strip():
                    original_text = para.text
                    encrypted_text = encrypt_text(para.text, secret)
                    if original_text != encrypted_text:  # 确认文本已被修改
                        para.text = encrypted_text
                        modified = True
                        print(f"Encrypted text: {encrypted_text}")
            
            if not modified:
                print("Warning: No text was encrypted")
                
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], 'encrypted_' + os.path.basename(file_path))
            doc.save(output_path)
            return output_path
        else:
            # 解密模式
            for para in doc.paragraphs:
                if para.text.strip():
                    print(f"Processing paragraph: {para.text}")  # 调试输出
                    # 检查段落中是否包含零宽字符
                    has_zero_width = any(c in para.text for c in [ZERO_WIDTH_SPACE, ZERO_WIDTH_JOINER])
                    if has_zero_width:
                        print(f"Found zero-width characters in paragraph")
                        try:
                            decrypted = decrypt_text(para.text)
                            if decrypted:
                                result.append(decrypted)
                                print(f"Successfully decrypted: {decrypted}")
                        except Exception as e:
                            print(f"Error decrypting paragraph: {str(e)}")
                            continue
            
            # 检查是否找到任何隐藏信息
            if not result:
                print("No hidden information found in any paragraph")
                return "未找到隐藏信息"
            
            final_result = '\n'.join(filter(None, result))
            print(f"Final decryption result: {final_result}")
            return final_result
            
    except Exception as e:
        print(f"Error processing DOCX: {str(e)}")
        raise

def process_pptx(file_path, mode='encrypt', secret=''):
    prs = Presentation(file_path)
    result = []
    
    if mode == 'encrypt':
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shape.text = encrypt_text(shape.text, secret)
    else:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    try:
                        decrypted = decrypt_text(shape.text)
                        if decrypted:
                            result.append(decrypted)
                    except:
                        continue
    
    if mode == 'encrypt':
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], 'encrypted_' + os.path.basename(file_path))
        prs.save(output_path)
        return output_path
    else:
        return '\n'.join(filter(None, result))

def process_pdf(file_path, mode='encrypt', secret=''):
    reader = PdfReader(file_path)
    result = []
    
    if mode == 'encrypt':
        writer = PdfWriter()
        for page in reader.pages:
            text = page.extract_text()
            # 在PDF中添加隐藏文本层
            page.merge_page(create_text_overlay(encrypt_text(text, secret)))
            writer.add_page(page)
        
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], 'encrypted_' + os.path.basename(file_path))
        with open(output_path, 'wb') as f:
            writer.write(f)
        return output_path
    else:
        for page in reader.pages:
            text = page.extract_text()
            try:
                decrypted = decrypt_text(text)
                if decrypted:
                    result.append(decrypted)
            except:
                continue
        return '\n'.join(result)

# 添加新的函数处理 .ppt 文件
def process_ppt(file_path, mode='encrypt', secret=''):
    # 创建临时文件路径
    temp_dir = tempfile.mkdtemp()
    temp_pptx = os.path.join(temp_dir, 'temp.pptx')
    
    try:
        # 使用 PowerPoint 应用程序转换 .ppt 为 .pptx
        powerpoint = win32com.client.Dispatch('PowerPoint.Application')
        ppt = powerpoint.Presentations.Open(file_path)
        ppt.SaveAs(temp_pptx, 24)  # 24 是 .pptx 格式的文件格式代码
        ppt.Close()
        powerpoint.Quit()
        
        # 使用已有的 process_pptx 函数处理转换后的文件
        if mode == 'encrypt':
            result = process_pptx(temp_pptx, mode, secret)
            
            # 如果是加密模式，需要将结果文件转回 .ppt 格式
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], 'encrypted_' + os.path.basename(file_path))
            powerpoint = win32com.client.Dispatch('PowerPoint.Application')
            ppt = powerpoint.Presentations.Open(result)
            ppt.SaveAs(output_path, 1)  # 1 是 .ppt 格式的文件格式代码
            ppt.Close()
            powerpoint.Quit()
            return output_path
        else:
            return process_pptx(temp_pptx, mode)
            
    finally:
        # 清理临时文件
        if os.path.exists(temp_pptx):
            os.remove(temp_pptx)
        if os.path.exists(temp_dir):
            os.rmdir(temp_dir)

# 添加新的路由
@app.route('/encrypt_file', methods=['POST'])
def encrypt_file():
    if 'file' not in request.files:
        return jsonify({'result': "未找到文件", 'action': '错误'})
    
    files = request.files.getlist('file')
    secret = request.form['secret']
    
    if not files or files[0].filename == '':
        return jsonify({'result': "未选择文件", 'action': '错误'})
    
    results = []
    for file in files:
        filename = file.filename
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        output_path = None
        
        try:
            file.save(filepath)
            
            if filename.endswith('.docx'):
                output_path = process_docx(filepath, 'encrypt', secret)
            elif filename.endswith('.pptx'):
                output_path = process_pptx(filepath, 'encrypt', secret)
            elif filename.endswith('.ppt'):
                output_path = process_ppt(filepath, 'encrypt', secret)
            elif filename.endswith('.pdf'):
                output_path = process_pdf(filepath, 'encrypt', secret)
            else:
                results.append(f"{filename}: 不支持的文件格式")
                continue
            
            # 读取加密后的文件
            with open(output_path, 'rb') as f:
                file_content = f.read()
            
            # 创建加密后的文件名
            encrypted_filename = f'encrypted_{filename}'
            
            # 保存加密后的文件到结果列表
            results.append({
                'filename': encrypted_filename,
                'content': file_content
            })
            
        except Exception as e:
            results.append(f"{filename}: 处理失败 - {str(e)}")
            
        finally:
            # 清理临时文件
            if filepath and os.path.exists(filepath):
                try:
                    os.remove(filepath)
                except:
                    pass
            if output_path and os.path.exists(output_path):
                try:
                    os.remove(output_path)
                except:
                    pass
    
    # 返回处理结果
    return jsonify({
        'result': results,
        'action': '加密完成'
    })

@app.route('/decrypt_file', methods=['POST'])
def decrypt_file():
    if 'file' not in request.files:
        return jsonify({'result': "未找到文件", 'action': '错误'})
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'result': "未选择文件", 'action': '错误'})
    
    filename = file.filename
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    
    try:
        # 保存上传的文件
        file.save(filepath)
        
        # 根据文件类型选择处理方法
        result = None
        if filename.endswith('.docx'):
            result = process_docx(filepath, 'decrypt')
        elif filename.endswith('.pptx'):
            result = process_pptx(filepath, 'decrypt')
        elif filename.endswith('.ppt'):
            result = process_ppt(filepath, 'decrypt')
        elif filename.endswith('.pdf'):
            result = process_pdf(filepath, 'decrypt')
        else:
            return jsonify({'result': "不支持的文件格式", 'action': '错误'})
        
        # 检查解密结果
        if not result:
            return jsonify({'result': "未找到隐藏信息", 'action': '文件解密'})
        
        print(f"Decryption result: {result}")  # 添加调试输出
        return jsonify({'result': result, 'action': '文件解密'})
        
    except Exception as e:
        print(f"Decryption error: {str(e)}")  # 添加错误日志
        return jsonify({'result': f"处理文件时出错: {str(e)}", 'action': '错误'})
        
    finally:
        # 清理临时文件
        try:
            if os.path.exists(filepath):
                os.remove(filepath)
        except Exception as e:
            print(f"Error cleaning up file: {str(e)}")

if __name__ == '__main__':
    app.run(debug=True) 