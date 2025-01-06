from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import BaseOxmlElement
from .base import WatermarkBase
import hashlib
import json
import os

class PPTWatermark(WatermarkBase):
    def __init__(self):
        super().__init__()
        self.watermark_locations = []

    def embed_watermark(self, file_path, watermark_data, password):
        """在PPT文档中嵌入水印"""
        try:
            prs = Presentation(file_path)
            
            # 准备要嵌入的核心数据
            core_data = {
                'content': watermark_data['content'],
                'user': watermark_data['user']
            }
            
            # 1. 在演示文稿属性中嵌入
            self._embed_in_properties(prs, core_data)
            
            # 2. 在幻灯片布局中嵌入
            self._embed_in_layouts(prs, core_data)
            
            # 3. 在演示文稿注释中嵌入
            self._embed_in_comments(prs, core_data)
            
            # 保存文档
            output_path = os.path.join(os.path.dirname(file_path), 
                                     f'processed_{os.path.basename(file_path)}')
            prs.save(output_path)
            return output_path
            
        except Exception as e:
            print(f"Error in embed_watermark: {str(e)}")
            raise

    def _embed_in_properties(self, prs, core_data):
        """在演示文稿属性中嵌入水印"""
        try:
            # 获取核心属性
            core_props = prs.core_properties
            
            # 在注释中嵌入数据
            core_props.comments = self._hide_in_property(json.dumps(core_data))
            
            print(f"Embedded in properties: {core_data}")  # 调试信息
            
        except Exception as e:
            print(f"Error embedding in properties: {str(e)}")

    def _embed_in_layouts(self, prs, core_data):
        """在幻灯片布局中嵌入水印"""
        try:
            # 遍历所有布局
            for layout in prs.slide_layouts:
                # 生成唯一标识
                layout_id = hashlib.sha256(
                    (json.dumps(core_data) + str(layout.name)).encode()
                ).hexdigest()[:8]
                
                # 在布局XML中嵌入信息
                layout._element.set('customData', 
                    self._hide_in_property(json.dumps(core_data)))
                layout._element.set('id', layout_id)
            
            print(f"Embedded in layouts: {core_data}")  # 调试信息
            
        except Exception as e:
            print(f"Error embedding in layouts: {str(e)}")

    def _embed_in_comments(self, prs, core_data):
        """在演示文稿注释中嵌入水印"""
        try:
            # 在第一张幻灯片添加隐藏注释
            if len(prs.slides) > 0:
                slide = prs.slides[0]
                comment = slide.notes_slide.notes_text_frame
                comment.text = self._hide_in_property(json.dumps(core_data))
            
            print(f"Embedded in comments: {core_data}")  # 调试信息
            
        except Exception as e:
            print(f"Error embedding in comments: {str(e)}")

    def extract_watermark(self, file_path):
        """从PPT文档中提取水印"""
        try:
            prs = Presentation(file_path)
            watermark_data = None
            
            # 1. 从属性中提取
            prop_data = self._extract_from_properties(prs)
            if prop_data:
                watermark_data = prop_data
            
            # 2. 从布局中提取
            layout_data = self._extract_from_layouts(prs)
            if layout_data:
                # 验证数据一致性
                if watermark_data and watermark_data != layout_data:
                    return None
                watermark_data = layout_data
            
            # 3. 从注释中提取
            comment_data = self._extract_from_comments(prs)
            if comment_data:
                # 验证数据一致性
                if watermark_data and watermark_data != comment_data:
                    return None
                watermark_data = comment_data
            
            return watermark_data
            
        except Exception as e:
            print(f"Error extracting watermark: {str(e)}")
            return None

    def _extract_from_properties(self, prs):
        """从演示文稿属性中提取水印"""
        try:
            if prs.core_properties.comments:
                hidden_data = self._reveal_from_property(prs.core_properties.comments)
                if hidden_data:
                    data = json.loads(hidden_data)
                    print(f"Extracted from properties: {data}")  # 调试信息
                    return data
            return None
        except Exception as e:
            print(f"Error extracting from properties: {str(e)}")
            return None

    def _extract_from_layouts(self, prs):
        """从布局中提取水印"""
        try:
            for layout in prs.slide_layouts:
                custom_data = layout._element.get('customData')
                if custom_data:
                    hidden_data = self._reveal_from_property(custom_data)
                    if hidden_data:
                        data = json.loads(hidden_data)
                        print(f"Extracted from layouts: {data}")  # 调试信息
                        return data
            return None
        except Exception as e:
            print(f"Error extracting from layouts: {str(e)}")
            return None

    def _extract_from_comments(self, prs):
        """从注释中提取水印"""
        try:
            if len(prs.slides) > 0:
                slide = prs.slides[0]
                if slide.has_notes_slide:
                    comment = slide.notes_slide.notes_text_frame.text
                    if comment:
                        hidden_data = self._reveal_from_property(comment)
                        if hidden_data:
                            data = json.loads(hidden_data)
                            print(f"Extracted from comments: {data}")  # 调试信息
                            return data
            return None
        except Exception as e:
            print(f"Error extracting from comments: {str(e)}")
            return None

    def _hide_in_property(self, text):
        """将信息隐藏在属性值中"""
        if not text:
            return ''
        # 使用特殊字符混淆数据
        chars = list(text)
        for i in range(len(chars)):
            if i % 2 == 0:
                chars.insert(i, '\u200b')
        return ''.join(chars)

    def _reveal_from_property(self, text):
        """从属性值中提取信息"""
        if not text:
            return ''
        return ''.join(c for c in text if c != '\u200b') 