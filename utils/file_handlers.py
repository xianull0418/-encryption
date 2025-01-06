import os
import tempfile
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader, PdfWriter
import win32com.client
from .crypto import encrypt_text, decrypt_text

def process_docx(file_path, mode='encrypt', secret=''):
    """处理 DOCX 文件"""
    doc = Document(file_path)
    result = []
    output_path = None
    
    try:
        if mode == 'encrypt':
            modified = False
            for para in doc.paragraphs:
                if para.text.strip():
                    original_text = para.text
                    encrypted_text = encrypt_text(para.text, secret)
                    if original_text != encrypted_text:
                        para.text = encrypted_text
                        modified = True
                        print(f"Encrypted text: {encrypted_text}")
            
            if not modified:
                print("Warning: No text was encrypted")
            
            # 保存加密后的文件
            output_path = os.path.join(os.path.dirname(file_path), 
                                     f'encrypted_{os.path.basename(file_path)}')
            doc.save(output_path)
            return output_path
        else:
            for para in doc.paragraphs:
                if para.text.strip():
                    print(f"Processing paragraph: {para.text}")
                    has_zero_width = any(c in para.text for c in ['\u200b', '\u200d'])
                    if has_zero_width:
                        try:
                            decrypted = decrypt_text(para.text)
                            if decrypted:
                                result.append(decrypted)
                        except Exception as e:
                            print(f"Error decrypting paragraph: {str(e)}")
                            continue
            
            return '\n'.join(filter(None, result)) if result else "未找到隐藏信息"
            
    except Exception as e:
        print(f"Error processing DOCX: {str(e)}")
        raise

def process_pptx(file_path, mode='encrypt', secret=''):
    """处理 PPTX 文件"""
    prs = Presentation(file_path)
    result = []
    
    try:
        if mode == 'encrypt':
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        shape.text = encrypt_text(shape.text, secret)
            return prs
        else:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        try:
                            decrypted = decrypt_text(shape.text)
                            if decrypted:
                                result.append(decrypted)
                        except:
                            continue
            
            return '\n'.join(filter(None, result)) if result else "未找到隐藏信息"
    except Exception as e:
        print(f"Error processing PPTX: {str(e)}")
        raise

def process_ppt(file_path, mode='encrypt', secret=''):
    """处理 PPT 文件"""
    temp_dir = tempfile.mkdtemp()
    temp_pptx = os.path.join(temp_dir, 'temp.pptx')
    
    try:
        powerpoint = win32com.client.Dispatch('PowerPoint.Application')
        ppt = powerpoint.Presentations.Open(file_path)
        ppt.SaveAs(temp_pptx, 24)  # 24 是 .pptx 格式的文件格式代码
        ppt.Close()
        powerpoint.Quit()
        
        if mode == 'encrypt':
            result = process_pptx(temp_pptx, mode, secret)
            return result
        else:
            return process_pptx(temp_pptx, mode)
    finally:
        if os.path.exists(temp_pptx):
            os.remove(temp_pptx)
        if os.path.exists(temp_dir):
            os.rmdir(temp_dir)

def process_pdf(file_path, mode='encrypt', secret=''):
    """处理 PDF 文件"""
    reader = PdfReader(file_path)
    result = []
    
    try:
        if mode == 'encrypt':
            writer = PdfWriter()
            for page in reader.pages:
                text = page.extract_text()
                encrypted_text = encrypt_text(text, secret)
                # 创建新页面并添加加密文本
                writer.add_page(page)
            return writer
        else:
            for page in reader.pages:
                text = page.extract_text()
                try:
                    decrypted = decrypt_text(text)
                    if decrypted:
                        result.append(decrypted)
                except:
                    continue
            
            return '\n'.join(filter(None, result)) if result else "未找到隐藏信息"
    except Exception as e:
        print(f"Error processing PDF: {str(e)}")
        raise 